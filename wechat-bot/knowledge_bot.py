import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import json
from docx import Document as DocxDocument
import pandas as pd
import pdfplumber
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings.base import Embeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
import requests
from dotenv import load_dotenv
import numpy as np
import time
from sentence_transformers import SentenceTransformer
from sklearn.feature_extraction.text import TfidfVectorizer
import jieba
from ssl_config import configure_ssl, get_ssl_session
import torch
import traceback
from embeddings import SmartEmbeddings, TfidfEmbeddings
from maxkb_client import MaxKBClient
from smart_kb import SmartKnowledgeBase

# 加载环境变量
load_dotenv()

# 设置环境变量以启用离线模式
os.environ['TRANSFORMERS_OFFLINE'] = '1'
os.environ['HF_DATASETS_OFFLINE'] = '1'
os.environ['HF_HUB_OFFLINE'] = '1'

class TfidfEmbeddings(Embeddings):
    """使用TF-IDF的文本嵌入实现（离线备选方案）"""
    
    def __init__(self, dimension: int = 1536):
        self.dimension = dimension
        self.vectorizer = TfidfVectorizer(
            max_features=dimension,
            token_pattern=r"(?u)\b\w+\b"
        )
        self.fitted = False
        
    def fit(self, texts: List[str]):
        """训练TF-IDF向量化器"""
        # 对中文文本进行分词
        processed_texts = []
        for text in texts:
            words = jieba.cut(text)
            processed_text = " ".join(words)
            processed_texts.append(processed_text)
            
        self.vectorizer.fit(processed_texts)
        self.fitted = True
        
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """生成文档嵌入向量"""
        if not self.fitted:
            self.fit(texts)
            
        # 处理输入文本
        processed_texts = []
        for text in texts:
            words = jieba.cut(text)
            processed_text = " ".join(words)
            processed_texts.append(processed_text)
            
        # 生成TF-IDF向量
        vectors = self.vectorizer.transform(processed_texts).toarray()
        
        # 填充或截断到指定维度
        padded_vectors = []
        for vector in vectors:
            if len(vector) < self.dimension:
                padded_vector = np.pad(vector, (0, self.dimension - len(vector)))
            else:
                padded_vector = vector[:self.dimension]
            padded_vectors.append(padded_vector.tolist())
            
        return padded_vectors
    
    def embed_query(self, text: str) -> List[float]:
        """生成查询文本的嵌入向量"""
        return self.embed_documents([text])[0]

class SmartEmbeddings(Embeddings):
    """智能文本嵌入实现，支持多种后备方案"""
    
    def __init__(self):
        print("初始化智能向量化方法...")
        self.model = None
        self.fallback = None
        self.initialize_embeddings()
        
    def initialize_embeddings(self):
        """初始化嵌入模型，包含多个备选方案"""
        # 设置模型缓存目录
        cache_dir = os.path.join(os.getcwd(), 'models')
        os.makedirs(cache_dir, exist_ok=True)
        
        # 模型加载顺序：本地缓存 -> 下载 -> TF-IDF备选
        model_candidates = [
            ('shibing624/text2vec-base-chinese', '中文模型'),
            ('sentence-transformers/all-MiniLM-L6-v2', '通用模型'),
            ('sentence-transformers/paraphrase-MiniLM-L3-v2', '轻量模型')
        ]
        
        for model_name, model_type in model_candidates:
            try:
                print(f"尝试加载{model_type} {model_name}...")
                self.model = SentenceTransformer(
                    model_name,
                    device='cpu' if not torch.cuda.is_available() else 'cuda',
                    cache_folder=cache_dir,
                    trust_remote_code=True
                )
                print(f"成功加载{model_type}")
                return
            except Exception as e:
                print(f"加载{model_type}失败: {str(e)}")
                continue
        
        # 如果所有模型都加载失败，使用TF-IDF作为备选方案
        print("所有模型加载失败，使用TF-IDF作为备选方案")
        self.fallback = TfidfEmbeddings()
        
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """生成文档嵌入向量"""
        try:
            if self.model is not None:
                embeddings = self.model.encode(texts, convert_to_tensor=False)
                return embeddings.tolist()
            else:
                return self.fallback.embed_documents(texts)
        except Exception as e:
            print(f"生成文档向量时出错: {e}")
            if self.fallback is None:
                self.fallback = TfidfEmbeddings()
            return self.fallback.embed_documents(texts)
    
    def embed_query(self, text: str) -> List[float]:
        """生成查询文本的嵌入向量"""
        return self.embed_documents([text])[0]

class SiliconFlowEmbeddings(Embeddings):
    """使用SiliconFlow API的文本嵌入实现"""
    
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.api_url = "https://api.siliconflow.cn/v1/chat/completions"
        self.session = get_ssl_session()
        
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """生成文档嵌入向量"""
        embeddings = []
        for text in texts:
            embedding = self.embed_query(text)
            embeddings.append(embedding)
        return embeddings
    
    def embed_query(self, text: str) -> List[float]:
        """生成查询文本的嵌入向量"""
        max_retries = 3
        current_retry = 0
        while current_retry < max_retries:
            try:
                response = self.session.post(
                    self.api_url,
                    headers={
                        "Content-Type": "application/json",
                        "Authorization": f"Bearer {self.api_key}"
                    },
                    json={
                        "model": "deepseek-ai/DeepSeek-R1-0528-Qwen3-8B",
                        "messages": [
                            {"role": "system", "content": "你是一个文本向量化助手。请将输入的文本转换为向量表示。"},
                            {"role": "user", "content": text}
                        ],
                        "stream": False,
                        "max_tokens": 2048,
                        "temperature": 0.7,
                        "top_p": 0.7,
                        "frequency_penalty": 0.5,
                        "presence_penalty": 0.0
                    },
                    timeout=60
                )
                
                if response.status_code == 200:
                    data = response.json()
                    content = data["choices"][0]["message"]["content"]
                    vector = [ord(c) / 255.0 for c in content[:1536]]
                    if len(vector) < 1536:
                        vector.extend([0.0] * (1536 - len(vector)))
                    else:
                        vector = vector[:1536]
                    return vector
                else:
                    print(f"API调用失败 (尝试 {current_retry + 1}/{max_retries}): HTTP {response.status_code}")
                    print(f"错误信息: {response.text}")
                    
            except requests.exceptions.Timeout:
                print(f"API请求超时 (尝试 {current_retry + 1}/{max_retries})")
            except requests.exceptions.RequestException as e:
                print(f"API请求异常 (尝试 {current_retry + 1}/{max_retries}): {e}")
            except Exception as e:
                print(f"生成嵌入向量时出错 (尝试 {current_retry + 1}/{max_retries}): {e}")
            
            current_retry += 1
            if current_retry < max_retries:
                print(f"等待 {current_retry * 2} 秒后重试...")
                time.sleep(current_retry * 2)
        
        print("所有重试都失败，返回随机向量")
        return list(np.random.rand(1536))

class TextLoader:
    """自定义的文本加载器"""
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        with open(self.file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        return [Document(page_content=text, metadata={"source": self.file_path})]

class DocxLoader:
    """自定义的Word文档加载器"""
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        doc = DocxDocument(self.file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return [Document(page_content=text, metadata={"source": self.file_path})]

class ExcelLoader:
    """Excel文档加载器"""
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        # 读取所有sheet
        dfs = pd.read_excel(self.file_path, sheet_name=None)
        texts = []
        
        # 处理每个sheet
        for sheet_name, df in dfs.items():
            # 将sheet名称添加为标题
            sheet_text = f"Sheet: {sheet_name}\n\n"
            
            # 添加列名
            sheet_text += "列: " + ", ".join(df.columns) + "\n\n"
            
            # 将数据转换为文本
            for _, row in df.iterrows():
                row_text = " | ".join([str(val) for val in row])
                sheet_text += row_text + "\n"
            
            texts.append(sheet_text)
        
        # 合并所有sheet的文本
        combined_text = "\n\n".join(texts)
        return [Document(page_content=combined_text, metadata={"source": self.file_path})]

class PDFLoader:
    """PDF文档加载器"""
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        texts = []
        with pdfplumber.open(self.file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    texts.append(text)
        
        combined_text = "\n\n".join(texts)
        return [Document(page_content=combined_text, metadata={"source": self.file_path})]

class PPTLoader:
    """PowerPoint文档加载器"""
    def __init__(self, file_path: str):
        self.file_path = file_path

    def load(self) -> List[Document]:
        prs = Presentation(self.file_path)
        texts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"Slide {i}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            texts.append(slide_text)
        
        combined_text = "\n\n".join(texts)
        return [Document(page_content=combined_text, metadata={"source": self.file_path})]

class KnowledgeBot:
    def __init__(self, knowledge_base_path: str = "knowledge_base"):
        """
        初始化知识库机器人
        
        Args:
            knowledge_base_path: 知识库路径
        """
        print("初始化知识库机器人...")
        
        # 初始化智能知识库
        try:
            self.kb = SmartKnowledgeBase(knowledge_base_path)
            print("知识库初始化完成")
        except Exception as e:
            print(f"知识库初始化失败: {str(e)}")
            print("将使用大模型作为备选方案")
            self.kb = None
            
    def query(self, question: str) -> str:
        """
        查询知识库并返回答案
        
        Args:
            question: 用户问题
            
        Returns:
            str: 答案
        """
        try:
            if not self.kb:
                print("\n=== 知识库不可用，直接使用大模型 ===")
                return self._get_llm_response(question)
                
            # 从知识库获取相关内容
            print("\n=== 正在查询知识库 ===")
            print(f"问题: {question}")
            
            # 使用更低的相似度阈值
            thresholds = [0.05, 0.03, 0.01]  # 从相对严格到宽松
            context = None
            best_threshold = None
            
            for threshold in thresholds:
                print(f"\n尝试相似度阈值: {threshold}")
                self.kb.vector_store.similarity_threshold = threshold
                try:
                    current_context = self.kb.query(question)
                    if current_context:
                        print(f"在阈值 {threshold} 下找到相关内容")
                        # 如果是第一次找到内容，或者当前内容更相关（基于长度），则更新
                        if not context or len(current_context) > len(context):
                            context = current_context
                            best_threshold = threshold
                except Exception as e:
                    print(f"在阈值 {threshold} 下查询失败: {str(e)}")
                    continue
            
            if not context:
                print("\n在所有阈值下都没有找到相关内容，使用大模型")
                return self._get_llm_response(question)
                
            print(f"\n使用最佳阈值 {best_threshold} 找到相关内容：")
            print(f"长度: {len(context)} 字符")
            print(f"内容预览: {context[:200]}...")
            
            # 直接返回知识库内容，不再使用大模型处理
            return context
            
        except Exception as e:
            print(f"\n查询知识库出错: {str(e)}")
            print("\n详细错误信息:")
            traceback.print_exc()
            print("\n发生错误，使用大模型作为备选方案")
            return self._get_llm_response(question)
            
    def add_document(self, file_path: str) -> bool:
        """
        添加文档到知识库
        
        Args:
            file_path: 文档路径
            
        Returns:
            bool: 是否添加成功
        """
        try:
            print(f"\n=== 开始添加文档到知识库 ===")
            print(f"文档路径: {file_path}")
            
            if not self.kb:
                print("错误: 知识库不可用")
                return False
            
            # 检查文件是否存在
            if not os.path.exists(file_path):
                print(f"错误: 文件不存在: {file_path}")
                return False
                
            # 检查文件大小
            file_size = os.path.getsize(file_path)
            print(f"文件大小: {file_size / 1024:.2f} KB")
            
            # 检查文件类型
            file_ext = Path(file_path).suffix.lower()
            print(f"文件类型: {file_ext}")
            
            # 尝试添加文档
            print("正在添加文档到知识库...")
            success = self.kb.add_document(file_path)
            
            if success:
                print("文档添加成功")
                # 获取更新后的统计信息
                stats = self.get_statistics()
                print("\n知识库统计信息:")
                print(f"- 文档总数: {stats['total_documents']}")
                print(f"- 文本块总数: {stats['total_chunks']}")
            else:
                print("文档添加失败")
            
            return success
            
        except Exception as e:
            print(f"\n添加文档时出错:")
            print(f"错误类型: {type(e).__name__}")
            print(f"错误信息: {str(e)}")
            print("\n详细错误信息:")
            traceback.print_exc()
            return False
            
    def _get_llm_response(self, question: str, context: str = None) -> str:
        """
        使用大模型生成回答
        
        Args:
            question: 用户问题
            context: 上下文信息（可选）
            
        Returns:
            str: 大模型生成的回答
        """
        try:
            print("\n=== 调用大模型生成回答 ===")
            # 构建提示词
            if context:
                print("使用知识库内容作为上下文")
                prompt = f"""请基于以下信息回答问题。如果信息不足以回答问题，请明确说明。

参考信息：
{context}

问题：{question}

请用简洁专业的语言回答。"""
            else:
                print("直接使用问题生成回答")
                prompt = question
                
            print(f"问题：{question}")
            
            # 调用大模型API
            url = "https://api.siliconflow.cn/v1/chat/completions"
            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {os.getenv('SILICONFLOW_API_KEY')}"
            }
            data = {
                "model": "deepseek-ai/DeepSeek-R1-0528-Qwen3-8B",
                "messages": [
                    {"role": "system", "content": "你是一个专业、友好的AI助手，请用简洁专业的语言回答问题。"},
                    {"role": "user", "content": prompt}
                ],
                "stream": False,
                "max_tokens": 2048,
                "temperature": 0.7,
                "top_p": 0.7,
                "frequency_penalty": 0.5,
                "presence_penalty": 0.0
            }
            
            print("\n发送请求到大模型...")
            response = requests.post(url, headers=headers, json=data)
            response.raise_for_status()
            result = response.json()
            
            if "choices" in result and len(result["choices"]) > 0:
                answer = result["choices"][0]["message"]["content"]
                print(f"\n大模型回答：\n{answer}")
                return answer
            else:
                print("\n错误：API返回的数据格式不正确")
                return "抱歉，我无法生成回答。"
                
        except Exception as e:
            print(f"\n调用大模型出错: {str(e)}")
            return f"抱歉，生成回答时出错: {str(e)}"

    def get_statistics(self) -> Dict[str, Any]:
        """
        获取知识库统计信息
        
        Returns:
            Dict[str, Any]: 包含统计信息的字典
        """
        try:
            if not self.kb:
                return {
                    "total_documents": 0,
                    "total_chunks": 0,
                    "documents": []
                }
            
            # 获取向量存储中的文档数量
            total_chunks = len(self.kb.vector_store.texts) if self.kb.vector_store else 0
            
            # 获取文档列表
            docs_dir = Path(self.kb.knowledge_base_path) / "docs"
            documents = []
            total_documents = 0
            
            if docs_dir.exists():
                for doc in docs_dir.glob("**/*"):
                    if doc.is_file():
                        total_documents += 1
                        doc_chunks = sum(1 for _ in self.kb.text_splitter.split_text(doc.read_text(encoding='utf-8')))
                        documents.append({
                            "path": str(doc.relative_to(docs_dir)),
                            "chunk_count": doc_chunks
                        })
            
            return {
                "total_documents": total_documents,
                "total_chunks": total_chunks,
                "documents": documents
            }
            
        except Exception as e:
            print(f"获取统计信息失败: {e}")
            return {
                "total_documents": 0,
                "total_chunks": 0,
                "documents": [],
                "error": str(e)
            }

# 测试代码
if __name__ == "__main__":
    try:
        # 创建知识库实例
        bot = KnowledgeBot()
        
        # 打印知识库信息
        print("\n=== 知识库信息 ===")
        stats = bot.get_statistics()
        print(json.dumps(stats, indent=2, ensure_ascii=False))
        
        # 添加文档示例
        print("\n=== 添加文档测试 ===")
        test_files = ["docs/test.txt"]
        results = bot.add_documents(test_files)
        for file_path, success in results.items():
            print(f"添加文档 {file_path}: {'成功' if success else '失败'}")
        
        # 搜索示例
        print("\n=== 搜索测试 ===")
        test_queries = [
            "什么是退换货政策？",
            "如何申请退货？",
            "运费说明是什么？"
        ]
        
        for query in test_queries:
            print(f"\n查询: {query}")
            # 搜索相似文档
            similar_docs = bot.search_similar(query)
            print("\n相似文档片段:")
            for doc in similar_docs:
                print(f"- 来源: {doc['source']}")
                print(f"  相似度: {doc['similarity']:.4f}")
                print(f"  内容: {doc['content'][:200]}...")
            
            # 生成回答
            answer = bot.query(query)
            print(f"\n回答: {answer}")
        
        # 更新后的知识库信息
        print("\n=== 更新后的知识库信息 ===")
        stats = bot.get_statistics()
        print(json.dumps(stats, indent=2, ensure_ascii=False))
        
    except Exception as e:
        print(f"\n程序运行出错: {e}")
        print("\n详细错误信息:")
        import traceback
        traceback.print_exc() 