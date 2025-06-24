print('当前加载的 document_processor.py 路径:', __file__)
import os
import pandas as pd
from typing import List, Dict, Any, Optional
from pptx import Presentation
from openpyxl import load_workbook
from docx import Document
from PyPDF2 import PdfReader
import magic
import logging
from dataclasses import dataclass
from tqdm import tqdm
from loguru import logger
import traceback
from ai_chunk_service import AIChunkService
import re

logger = logging.getLogger(__name__)

@dataclass
class DocumentChunk:
    """文档分块类，用于存储文档内容和元数据"""
    content: str  # 文档内容
    metadata: Dict[str, Any]  # 元数据（如来源、类型等）

class DocumentProcessor:
    """文档处理器，支持多种文件格式的读取和处理"""
    
    SUPPORTED_EXTENSIONS = {
        # 文本文件
        '.txt': 'text/plain',
        # Microsoft Office
        '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        '.xls': 'application/vnd.ms-excel',
        '.ppt': 'application/vnd.ms-powerpoint',
        '.doc': 'application/msword',
        # PDF
        '.pdf': 'application/pdf',
        # 其他文本格式
        '.md': 'text/markdown',
        '.csv': 'text/csv',
        '.json': 'application/json',
    }

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        """
        初始化文档处理器
        :param chunk_size: 文本分块大小
        :param chunk_overlap: 分块重叠大小
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def process_file(self, file_path: str) -> List[DocumentChunk]:
        """
        处理单个文件
        :param file_path: 文件路径
        :return: 文档分块列表
        """
        try:
            # 检测文件类型
            file_type = magic.from_file(file_path, mime=True)
            extension = os.path.splitext(file_path)[1].lower()

            # 验证文件类型
            if extension not in self.SUPPORTED_EXTENSIONS:
                raise ValueError(f"Unsupported file extension: {extension}")

            # 根据文件类型调用相应的处理方法
            if extension in ['.xlsx', '.xls']:
                return self._process_excel(file_path)
            elif extension in ['.pptx', '.ppt']:
                return self._process_powerpoint(file_path)
            elif extension in ['.docx', '.doc']:
                return self._process_word(file_path)
            elif extension == '.pdf':
                return self._process_pdf(file_path)
            elif extension == '.csv':
                return self._process_csv(file_path)
            else:  # 处理纯文本文件
                return self._process_text(file_path)

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise

    def _process_excel(self, file_path: str) -> List[DocumentChunk]:
        """处理Excel文件"""
        chunks = []
        try:
            logger.info(f"开始处理Excel文件: {file_path}")
            
            # 使用pandas读取Excel文件
            logger.info("读取Excel文件...")
            excel_file = pd.ExcelFile(file_path)
            logger.info(f"发现 {len(excel_file.sheet_names)} 个工作表: {excel_file.sheet_names}")
            
            # 处理每个工作表
            for sheet_name in excel_file.sheet_names:
                logger.info(f"\n处理工作表: {sheet_name}")
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 获取表头
                headers = df.columns.tolist()
                logger.info(f"列名: {headers}")
                logger.info(f"数据行数: {len(df)}")
                
                # 处理每一行
                for index, row in df.iterrows():
                    # 构建行文本
                    row_text = " | ".join([f"{headers[i]}: {str(row[i])}" for i in range(len(headers))])
                    
                    metadata = {
                        'source': file_path,
                        'sheet_name': sheet_name,
                        'row_index': index,
                        'type': 'excel'
                    }
                    
                    logger.debug(f"行 {index + 1}:")
                    logger.debug(f"- 文本: {row_text}")
                    logger.debug(f"- 元数据: {metadata}")
                    
                    chunks.append(DocumentChunk(content=row_text, metadata=metadata))
                
                # 添加表格概述
                summary = f"工作表 '{sheet_name}' 包含 {len(df)} 行数据，列名为：{', '.join(headers)}"
                metadata = {
                    'source': file_path,
                    'sheet_name': sheet_name,
                    'type': 'excel_summary'
                }
                
                logger.info(f"添加表格概述: {summary}")
                chunks.append(DocumentChunk(content=summary, metadata=metadata))
                
            logger.info(f"Excel处理完成，共生成 {len(chunks)} 个文本块")
            
            # 记录一些示例块
            if chunks:
                logger.info("\n文本块示例:")
                for i, chunk in enumerate(chunks[:3]):  # 只显示前3个
                    logger.info(f"块 {i+1}:")
                    logger.info(f"- 内容: {chunk.content}")
                    logger.info(f"- 元数据: {chunk.metadata}")
            
        except Exception as e:
            logger.error(f"处理Excel文件时出错: {str(e)}")
            logger.error(traceback.format_exc())
            raise
            
        return chunks

    def _process_powerpoint(self, file_path: str) -> List[DocumentChunk]:
        """处理PowerPoint文件"""
        chunks = []
        try:
            presentation = Presentation(file_path)
            
            for slide_number, slide in enumerate(presentation.slides, 1):
                slide_text = []
                
                # 处理标题
                if slide.shapes.title:
                    slide_text.append(f"标题: {slide.shapes.title.text}")
                
                # 处理其他形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    content = "\n".join(slide_text)
                    metadata = {
                        'source': file_path,
                        'slide_number': slide_number,
                        'type': 'powerpoint'
                    }
                    chunks.append(DocumentChunk(content=content, metadata=metadata))
                
            # 添加PPT概述
            summary = f"演示文稿包含 {len(presentation.slides)} 张幻灯片"
            metadata = {
                'source': file_path,
                'type': 'powerpoint_summary'
            }
            chunks.append(DocumentChunk(content=summary, metadata=metadata))
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
            raise
            
        return chunks

    def _process_word(self, file_path: str) -> List[DocumentChunk]:
        """处理Word文件"""
        chunks = []
        try:
            # 根据文件扩展名选择处理方法
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext == '.docx':
                doc = Document(file_path)
                content = ""
                for para in doc.paragraphs:
                    content += para.text + "\n"
            elif file_ext == '.doc':
                # 尝试多种方法读取.doc文件
                content = self._read_doc_file(file_path)
            else:
                raise ValueError(f"Unsupported file extension: {file_ext}")
            
            if not content.strip():
                logger.warning(f"Word文件内容为空: {file_path}")
                return []
            
            # 按段落分割并创建块
            current_chunk = []
            current_length = 0
            
            paragraphs = content.split('\n')
            for para in paragraphs:
                text = para.strip()
                if not text:
                    continue
                    
                # 如果当前块加上新段落超过块大小，创建新块
                if current_length + len(text) > self.chunk_size:
                    if current_chunk:
                        content = "\n".join(current_chunk)
                        metadata = {
                            'source': file_path,
                            'type': 'word'
                        }
                        chunks.append(DocumentChunk(content=content, metadata=metadata))
                    current_chunk = [text]
                    current_length = len(text)
                else:
                    current_chunk.append(text)
                    current_length += len(text)
            
            # 处理最后一个块
            if current_chunk:
                content = "\n".join(current_chunk)
                metadata = {
                    'source': file_path,
                    'type': 'word'
                }
                chunks.append(DocumentChunk(content=content, metadata=metadata))
                
        except Exception as e:
            logger.error(f"Error processing Word file {file_path}: {str(e)}")
            raise
            
        return chunks
    
    def _read_doc_file(self, file_path: str) -> str:
        """读取.doc文件"""
        try:
            # 首先尝试使用python-docx2txt库
            try:
                import docx2txt
                content = docx2txt.process(file_path)
                if content and content.strip():
                    logger.info("使用docx2txt成功读取.doc文件")
                    return content
            except Exception as e:
                logger.warning(f"docx2txt读取.doc文件失败: {str(e)}")
            
            # 尝试使用olefile读取.doc文件
            try:
                import olefile
                if olefile.isOleFile(file_path):
                    ole = olefile.OleFileIO(file_path)
                    # 尝试读取WordDocument流
                    if ole.exists('WordDocument'):
                        word_doc = ole.openstream('WordDocument').read()
                        # 尝试读取内容
                        content = self._extract_text_from_ole(ole)
                        if content and content.strip():
                            logger.info("使用olefile成功读取.doc文件")
                            return content
                    ole.close()
            except ImportError:
                logger.warning("olefile库未安装")
            except Exception as e:
                logger.warning(f"olefile读取.doc文件失败: {str(e)}")
            
            # 尝试使用mammoth库
            try:
                import mammoth
                with open(file_path, "rb") as docx_file:
                    result = mammoth.extract_raw_text(docx_file)
                    if result.value and result.value.strip():
                        logger.info("使用mammoth成功读取.doc文件")
                        return result.value
            except ImportError:
                logger.warning("mammoth库未安装")
            except Exception as e:
                logger.warning(f"mammoth读取.doc文件失败: {str(e)}")
            
            # 尝试使用antiword（Linux/Unix工具）
            try:
                import subprocess
                result = subprocess.run(['antiword', file_path], 
                                      capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    logger.info("使用antiword成功读取.doc文件")
                    return result.stdout
            except (ImportError, FileNotFoundError, subprocess.TimeoutExpired) as e:
                logger.warning(f"antiword未安装或执行失败: {str(e)}")
            
            # 尝试使用catdoc（Linux/Unix工具）
            try:
                import subprocess
                result = subprocess.run(['catdoc', file_path], 
                                      capture_output=True, text=True, timeout=30)
                if result.returncode == 0 and result.stdout.strip():
                    logger.info("使用catdoc成功读取.doc文件")
                    return result.stdout
            except (ImportError, FileNotFoundError, subprocess.TimeoutExpired) as e:
                logger.warning(f"catdoc未安装或执行失败: {str(e)}")
            
            # 如果所有方法都失败，抛出详细的异常信息
            error_msg = (
                "无法读取.doc文件。请尝试以下解决方案：\n"
                "1. 将.doc文件转换为.docx格式（推荐）：\n"
                "   - 用Microsoft Word打开文件，另存为.docx格式\n"
                "   - 或使用在线转换工具：https://convertio.co/doc-docx/\n"
                "2. 在Windows系统上安装LibreOffice，然后使用命令行转换：\n"
                "   soffice --headless --convert-to docx:MS Word 2007 XML 文件路径.doc\n"
                "3. 使用WPS Office或其他支持.doc的办公软件打开并另存为.docx\n"
                "4. 如果文件很重要，建议使用Microsoft Word进行转换"
            )
            raise Exception(error_msg)
            
        except Exception as e:
            logger.error(f"读取.doc文件失败: {str(e)}")
            raise
    
    def _extract_text_from_ole(self, ole):
        """从OLE文件中提取文本"""
        try:
            content = ""
            
            # 尝试读取不同的流
            streams_to_try = ['WordDocument', 'Contents', '1Table', '0Table']
            
            for stream_name in streams_to_try:
                if ole.exists(stream_name):
                    try:
                        stream = ole.openstream(stream_name)
                        data = stream.read()
                        # 尝试解码为文本
                        try:
                            text = data.decode('utf-8', errors='ignore')
                            # 清理文本
                            text = self._clean_ole_text(text)
                            if text.strip():
                                content += text + "\n"
                        except:
                            # 如果UTF-8失败，尝试其他编码
                            try:
                                text = data.decode('gbk', errors='ignore')
                                text = self._clean_ole_text(text)
                                if text.strip():
                                    content += text + "\n"
                            except:
                                pass
                        stream.close()
                    except:
                        continue
            
            return content
            
        except Exception as e:
            logger.warning(f"从OLE文件提取文本失败: {str(e)}")
            return ""
    
    def _clean_ole_text(self, text):
        """清理从OLE文件提取的文本"""
        # 移除控制字符
        import re
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
        # 移除多余的空白字符
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _process_pdf(self, file_path: str) -> List[DocumentChunk]:
        """处理PDF文件（全文件合并，AI分块，特殊字符清理，循环调用AI）"""
        chunks = []
        try:
            print('!!! PDF DEBUG ENTRY !!!')
            reader = PdfReader(file_path)
            ai_service = AIChunkService()
            max_ai_chars = int(3000)
            # 1. 合并所有页文本
            all_text = ''
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    all_text += text + '\n'
            # 2. 清理特殊字符
            all_text = re.sub(r'[\u2000-\u200F\u2028-\u202F\u205F-\u206F\uFEFF\x00-\x1F]', '', all_text)
            all_text = all_text.strip()
            logger.info(f"all_text type: {type(all_text)}, len: {len(all_text)}")
            if not isinstance(all_text, str):
                raise TypeError(f"all_text is not str, but {type(all_text)}")
            if not all_text:
                return []
            # 3. 按最大长度分块循环调用AI
            start = 0
            while start < len(all_text):
                start = int(start)
                end = int(start + int(max_ai_chars))
                logger.info(f"PDF分块循环: start={start}, end={end}, max_ai_chars={max_ai_chars}, type(start)={type(start)}, type(end)={type(end)}")
                assert isinstance(start, int) and isinstance(end, int), f'start={start}, end={end}, type(start)={type(start)}, type(end)={type(end)}'
                sub_text = all_text[start:end]
                # 保证不截断句子，向后扩展到最近的句号
                if end < len(all_text):
                    last_punc = max(sub_text.rfind('。'), sub_text.rfind('.'), sub_text.rfind('\n'))
                    if last_punc > 0 and last_punc > max_ai_chars * 0.5:
                        sub_text = sub_text[:last_punc+1]
                # 4. 调用AI分块
                ai_chunks = ai_service.chunk_with_ai(sub_text)
                if ai_chunks:
                    for ai_chunk in ai_chunks:
                        content = ai_chunk.get('content', '').strip()
                        if content:
                            metadata = {
                                'source': file_path,
                                'type': 'pdf_ai',
                            }
                            chunks.append(DocumentChunk(content=content, metadata=metadata))
                # 防止死循环：如果分块长度为0，强制跳步
                if len(sub_text) == 0:
                    start += max_ai_chars
                else:
                    start += len(sub_text)
        except Exception as e:
            logger.error(f"Error processing PDF file {file_path}: {str(e)}")
            logger.error(traceback.format_exc())
            raise
        return chunks

    def _process_csv(self, file_path: str) -> List[DocumentChunk]:
        """处理CSV文件"""
        chunks = []
        try:
            df = pd.read_csv(file_path)
            headers = df.columns.tolist()
            
            # 处理每一行
            for index, row in df.iterrows():
                # 构建行文本
                row_text = " | ".join([f"{headers[i]}: {str(row[i])}" for i in range(len(headers))])
                
                metadata = {
                    'source': file_path,
                    'row_index': index,
                    'type': 'csv'
                }
                
                chunks.append(DocumentChunk(content=row_text, metadata=metadata))
            
            # 添加CSV概述
            summary = f"CSV文件包含 {len(df)} 行数据，列名为：{', '.join(headers)}"
            metadata = {
                'source': file_path,
                'type': 'csv_summary'
            }
            chunks.append(DocumentChunk(content=summary, metadata=metadata))
            
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {str(e)}")
            raise
            
        return chunks

    def _process_text(self, file_path: str) -> List[DocumentChunk]:
        """处理文本文件"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
                
            # 按段落分割
            paragraphs = text.split('\n\n')
            current_chunk = []
            current_length = 0
            
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue
                    
                # 如果当前块加上新段落超过块大小，创建新块
                if current_length + len(para) > self.chunk_size:
                    if current_chunk:
                        content = "\n".join(current_chunk)
                        metadata = {
                            'source': file_path,
                            'type': 'text'
                        }
                        chunks.append(DocumentChunk(content=content, metadata=metadata))
                    current_chunk = [para]
                    current_length = len(para)
                else:
                    current_chunk.append(para)
                    current_length += len(para)
            
            # 处理最后一个块
            if current_chunk:
                content = "\n".join(current_chunk)
                metadata = {
                    'source': file_path,
                    'type': 'text'
                }
                chunks.append(DocumentChunk(content=content, metadata=metadata))
                
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {str(e)}")
            raise
            
        return chunks

    def process_directory(self, directory_path: str) -> List[DocumentChunk]:
        """
        处理目录下的所有支持的文件
        :param directory_path: 目录路径
        :return: 所有文档分块的列表
        """
        all_chunks = []
        
        # 获取所有文件
        files = []
        for root, _, filenames in os.walk(directory_path):
            for filename in filenames:
                ext = os.path.splitext(filename)[1].lower()
                if ext in self.SUPPORTED_EXTENSIONS:
                    files.append(os.path.join(root, filename))
        
        # 使用tqdm显示处理进度
        for file_path in tqdm(files, desc="Processing documents"):
            try:
                chunks = self.process_file(file_path)
                all_chunks.extend(chunks)
            except Exception as e:
                logger.error(f"Error processing {file_path}: {str(e)}")
                continue
        
        return all_chunks 